#!/usr/bin/env python3
"""Build MM Calendar deck — unified template, 8 slides, aspect-safe images."""
from pathlib import Path
from typing import List, Optional, Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
ASSETS = REPO / "mm_calendar" / "documentation" / "presentation_assets"
OUT = REPO / "mm_calendar" / "documentation" / "MM_Calendar_Tool_Presentation.pptx"
DESKTOP_DIR = Path.home() / "Desktop" / "MM_Calendar_Presentation"
DESKTOP_PPTX = Path.home() / "Desktop" / "MM_Calendar_Tool_Presentation.pptx"

BG = RGBColor(0xF7, 0xF5, 0xFB)
DARK = RGBColor(0x1E, 0x1B, 0x2E)
MUTED = RGBColor(0x5C, 0x57, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PLAN = RGBColor(0x5B, 0x21, 0xB6)
MONDAY = RGBColor(0x0E, 0x94, 0x88)
SC = RGBColor(0xDB, 0x27, 0x77)
INPUT = RGBColor(0xB4, 0x53, 0x09)  # amber — economy inputs

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.4)
HEADER_H = Inches(1.0)

prs: Presentation


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=BG) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def slide_header(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, HEADER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(MARGIN, Inches(0.16), SLIDE_W - MARGIN * 2, Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(MARGIN, Inches(0.62), SLIDE_W - MARGIN * 2, Inches(0.32))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = RGBColor(0xD1, 0xD5, 0xDB)


def label_strip(slide, top: Emu, text: str, rgb: RGBColor) -> Emu:
    label_h = Inches(0.32)
    box_w = SLIDE_W - MARGIN * 2
    lab = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, top, box_w, label_h)
    lab.fill.solid()
    lab.fill.fore_color.rgb = rgb
    lab.line.fill.background()
    lt = slide.shapes.add_textbox(MARGIN + Inches(0.12), top + Inches(0.05), box_w, label_h)
    lp = lt.text_frame.paragraphs[0]
    lp.text = text.upper()
    lp.font.size = Pt(10)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    return top + label_h + Inches(0.06)


def image_aspect(path: Path) -> float:
    with Image.open(path) as im:
        w, h = im.size
    return w / h if h else 1.0


def fit_picture(slide, path: Path, box_left, box_top, box_w, box_h) -> None:
    asp = image_aspect(path)
    box_w_f = box_w / Emu(914400)
    box_h_f = box_h / Emu(914400)
    if box_w_f / box_h_f >= asp:
        h_in = box_h_f
        w_in = h_in * asp
    else:
        w_in = box_w_f
        h_in = w_in / asp
    left_in = box_left / Emu(914400) + (box_w_f - w_in) / 2
    top_in = box_top / Emu(914400) + (box_h_f - h_in) / 2
    slide.shapes.add_picture(str(path), Inches(left_in), Inches(top_in), width=Inches(w_in))


def add_image_slide(
    title: str,
    subtitle: str,
    image_path: Path,
    strip_label: str,
    strip_rgb: RGBColor,
    caption: str,
) -> None:
    slide = blank_slide()
    fill_bg(slide)
    slide_header(slide, title, subtitle)
    if not image_path.is_file():
        return
    caption_h = Inches(0.45)
    img_top = label_strip(slide, HEADER_H + Inches(0.08), strip_label, strip_rgb)
    img_bottom = SLIDE_H - MARGIN - caption_h - Inches(0.08)
    box_h = img_bottom - img_top
    box_w = SLIDE_W - MARGIN * 2
    fit_picture(slide, image_path, MARGIN, img_top, box_w, box_h)
    cap = slide.shapes.add_textbox(MARGIN, SLIDE_H - MARGIN - caption_h, SLIDE_W - MARGIN * 2, caption_h)
    cp = cap.text_frame.paragraphs[0]
    cp.text = caption
    cp.font.size = Pt(12)
    cp.font.color.rgb = MUTED


def add_dual_image_slide(
    title: str,
    subtitle: str,
    left_path: Path,
    right_path: Path,
    strip_label: str,
) -> None:
    slide = blank_slide()
    fill_bg(slide)
    slide_header(slide, title, subtitle)
    img_top = label_strip(slide, HEADER_H + Inches(0.08), strip_label, INPUT)
    img_bottom = SLIDE_H - MARGIN - Inches(0.08)
    box_h = img_bottom - img_top
    gap = Inches(0.25)
    col_w = (SLIDE_W - MARGIN * 2 - gap) / 2
    for path, col in ((left_path, 0), (right_path, 1)):
        if path.is_file():
            left = MARGIN + col * (col_w + gap)
            fit_picture(slide, path, left, img_top, col_w, box_h)


def add_highlights_grid_slide(title: str, subtitle: str, paths: Sequence[Path]) -> None:
    slide = blank_slide()
    fill_bg(slide)
    slide_header(slide, title, subtitle)
    img_top = label_strip(slide, HEADER_H + Inches(0.08), "Monthly guideline inputs", INPUT)
    img_bottom = SLIDE_H - MARGIN - Inches(0.08)
    box_h = img_bottom - img_top
    cols, rows = 5, 2
    gap = Inches(0.12)
    box_w = SLIDE_W - MARGIN * 2
    cell_w = (box_w - gap * (cols - 1)) / cols
    cell_h = (box_h - gap * (rows - 1)) / rows
    for idx, path in enumerate(paths[: cols * rows]):
        if not path.is_file():
            continue
        r, c = divmod(idx, cols)
        left = MARGIN + c * (cell_w + gap)
        top = img_top + r * (cell_h + gap)
        fit_picture(slide, path, left, top, cell_w, cell_h)


def add_cover() -> None:
    s1 = blank_slide()
    fill_bg(s1, DARK)
    band = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.9), SLIDE_W, Inches(1.85))
    band.fill.solid()
    band.fill.fore_color.rgb = PLAN
    band.line.fill.background()
    for y, text, size, color, bold in [
        (1.4, "MM Calendar Tool", 42, WHITE, True),
        (3.05, "Plan  →  Monday  →  Smart Calendar", 24, WHITE, False),
        (5.2, "Rules · team handoff · LiveOps — kept in sync", 15, RGBColor(0xC4, 0xB5, 0xFD), False),
    ]:
        t = s1.shapes.add_textbox(Inches(0.6), Inches(y), Inches(12.1), Inches(0.9))
        p = t.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER


def add_pipeline_slide() -> None:
    slide = blank_slide()
    fill_bg(slide)
    slide_header(
        slide,
        "Three surfaces — same pipeline",
        "Inputs → builder → Monday → Smart Calendar (full-width screenshots follow).",
    )
    items = [
        ("0 · Economy & monthly inputs", INPUT, "Guidelines & highlights feed the builder"),
        ("1 · Planning dashboard", PLAN, "Builder output · rules & validation"),
        ("2 · Monday MM calendar", MONDAY, "Team handoff · Description · Add to SC"),
        ("3 · Smart Calendar", SC, "LiveOps config · verified vs Monday"),
    ]
    y = Inches(1.22)
    for title, col, sub in items:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, y, SLIDE_W - MARGIN * 2, Inches(0.78)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        t = slide.shapes.add_textbox(MARGIN + Inches(0.2), y + Inches(0.1), Inches(8), Inches(0.32))
        tp = t.text_frame.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(18)
        tp.font.bold = True
        tp.font.color.rgb = WHITE
        s = slide.shapes.add_textbox(MARGIN + Inches(0.2), y + Inches(0.42), Inches(10), Inches(0.28))
        sp = s.text_frame.paragraphs[0]
        sp.text = sub
        sp.font.size = Pt(13)
        sp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        y += Inches(0.92)


def add_loop_slide() -> None:
    slide = blank_slide()
    fill_bg(slide)
    slide_header(slide, "Weekly loop", "Cursor agents help build, fill Monday, and verify")
    steps = [
        ("Economy", "Caps & prize bank"),
        ("Build", "Dashboard + validation"),
        ("Monday", "Finalize rows · Add to SC"),
        ("LiveOps", "Smart Calendar config"),
        ("Verify", "Monday ↔ SC per day"),
    ]
    y = Inches(1.25)
    for n, (who, what) in enumerate(steps, 1):
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, y, Inches(0.38), Inches(0.38))
        circ.fill.solid()
        circ.fill.fore_color.rgb = PLAN
        circ.line.fill.background()
        nt = slide.shapes.add_textbox(MARGIN + Inches(0.05), y + Inches(0.04), Inches(0.32), Inches(0.32))
        np = nt.text_frame.paragraphs[0]
        np.text = str(n)
        np.font.size = Pt(13)
        np.font.bold = True
        np.font.color.rgb = WHITE
        wt = slide.shapes.add_textbox(Inches(1.05), y, Inches(2.0), Inches(0.4))
        wt.text_frame.paragraphs[0].text = who
        wt.text_frame.paragraphs[0].font.size = Pt(15)
        wt.text_frame.paragraphs[0].font.bold = True
        dt = slide.shapes.add_textbox(Inches(3.1), y, Inches(9.5), Inches(0.4))
        dt.text_frame.paragraphs[0].text = what
        dt.text_frame.paragraphs[0].font.size = Pt(15)
        dt.text_frame.paragraphs[0].font.color.rgb = MUTED
        y += Inches(0.62)


def highlight_paths() -> List[Path]:
    return sorted(ASSETS.glob("monthly_highlight_*.png"))


def main() -> None:
    global prs
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_cover()
    add_pipeline_slide()
    add_dual_image_slide(
        "Economy guidelines",
        "1st input of the tool",
        ASSETS / "economy_guidelines_1.png",
        ASSETS / "economy_guidelines_2.png",
        "Nivi monthly guidelines · caps & card bank",
    )
    add_highlights_grid_slide(
        "Monthly highlights",
        "2nd input of the tool",
        highlight_paths(),
    )
    add_image_slide(
        "Planning dashboard",
        "Output of the builder after validation",
        ASSETS / "mm_dashboard.png",
        "Repo dashboard · August 2026",
        PLAN,
        "Seasons, promos, and notes — encoded from the rule library.",
    )
    add_image_slide(
        "Monday — MM calendar",
        "Swimlanes by product lane; one row per promo",
        ASSETS / "monday_swimlanes.png",
        "Monday.com · week 9–15 Aug",
        MONDAY,
        "Description holds prizes & SKUs; economists review before Add to SC.",
    )
    add_image_slide(
        "Smart Calendar",
        "What LiveOps configured in the game (native aspect ratio)",
        ASSETS / "smart_calendar_week.png",
        "Smart Calendar · week 3–9 Aug",
        SC,
        "Automated check: promo names and Description text vs Monday.",
    )
    add_loop_slide()

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    prs.save(str(DESKTOP_PPTX))
    DESKTOP_DIR.mkdir(parents=True, exist_ok=True)
    print(f"Saved {len(prs.slides)} slides → {OUT}")


if __name__ == "__main__":
    main()
